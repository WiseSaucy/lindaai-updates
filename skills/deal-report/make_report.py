#!/usr/bin/env python3
"""
LindaAI — RV Park deal report generator.

Reads a filled RV_Park_Underwriting workbook, RECOMPUTES every metric in Python
(openpyxl does not evaluate Excel formulas), and renders Linda-style reports:

  • One-page PDF  (reportlab)  — offer summary, 5 key metrics, snapshot,
    normalized vs seller, offer comparison, red flags, value-add.
  • Slide deck    (python-pptx) — ~6 slides of the same.

Usage:
  python3 make_report.py "Deal - Example RV Park.xlsx" --pdf --pptx
  # outputs: "Deal - Example RV Park.pdf" and/or "Deal - Example RV Park.pptx"

Requires:  pip install -r requirements.txt   (openpyxl, reportlab, python-pptx)
"""
import argparse
import os
import sys

try:
    import openpyxl
except ImportError:
    sys.exit("Missing dependency: openpyxl (pip install -r requirements.txt)")


# --------------------------------------------------------------------------- #
# finance helpers
# --------------------------------------------------------------------------- #
def pmt(rate, nper, pv):
    if nper == 0:
        return 0.0
    if rate == 0:
        return -pv / nper
    return -(pv * rate) / (1 - (1 + rate) ** (-nper))


def irr(cashflows, lo=-0.95, hi=1.0, it=200):
    def npv(r):
        return sum(c / (1 + r) ** t for t, c in enumerate(cashflows))
    if npv(lo) * npv(hi) > 0:
        return None
    for _ in range(it):
        mid = (lo + hi) / 2
        if npv(lo) * npv(mid) <= 0:
            hi = mid
        else:
            lo = mid
    return (lo + hi) / 2


def gv(ws, cell, default=0.0):
    v = ws[cell].value
    if v is None:
        return default
    if isinstance(v, str):
        try:
            return float(v.replace(",", "").replace("$", "").replace("%", ""))
        except ValueError:
            return default
    return v


def money(x):
    try:
        return f"${x:,.0f}"
    except (TypeError, ValueError):
        return "—"


def pct(x, dp=1):
    try:
        return f"{x*100:.{dp}f}%"
    except (TypeError, ValueError):
        return "—"


def mult(x):
    try:
        return f"{x:.2f}x"
    except (TypeError, ValueError):
        return "—"


# --------------------------------------------------------------------------- #
# recompute the whole model
# --------------------------------------------------------------------------- #
def compute(wb):
    uw = wb["Underwriting"]
    d = {}
    d["name"] = uw["B3"].value or "RV Park"
    d["loc"] = uw["B4"].value or ""
    sites = gv(uw, "B8", 1) or 1
    price = gv(uw, "B9")
    rent = gv(uw, "B11")
    vac = gv(uw, "B12")
    other_m = gv(uw, "B13")
    mgmt_pct = gv(uw, "B14")
    res_site = gv(uw, "B15")
    rent_g = gv(uw, "B16")
    exp_g = gv(uw, "B17")
    use_detail = str(uw["B18"].value or "N").strip().upper()

    gpr = sites * rent * 12
    if use_detail == "Y" and "Income Detail" in wb.sheetnames:
        inc = wb["Income Detail"]
        lt = sum(gv(inc, f"B{r}") * gv(inc, f"C{r}") * 12 * gv(inc, f"D{r}") for r in (6, 7))
        tsites, nightly = gv(inc, "B11"), gv(inc, "B12")
        trans = sum(tsites * nightly * gv(inc, f"B{r}") * gv(inc, f"C{r}") for r in range(15, 27))
        gpr = lt + trans

    other_a = other_m * 12
    egi = gpr * (1 - vac) + other_a
    tax = gv(uw, "B27"); ins = gv(uw, "B28"); util = gv(uw, "B29")
    rm = gv(uw, "B30"); pay = gv(uw, "B31"); admin = gv(uw, "B32")
    mgmt = mgmt_pct * egi
    reserves = res_site * sites
    opex = tax + ins + util + rm + pay + admin + mgmt + reserves
    noi = egi - opex

    ltv = gv(uw, "B44"); intr = gv(uw, "B45"); amort = gv(uw, "B46")
    close_pct = gv(uw, "B47"); capex = gv(uw, "B48")
    hold = int(gv(uw, "B49", 5) or 5); exitcap = gv(uw, "B50"); sellcost = gv(uw, "B51")
    loan = price * ltv; down = price - loan; closing = price * close_pct
    cash = down + closing + capex
    pmt_m = pmt(intr / 12, amort * 12, -loan); ads = pmt_m * 12

    d.update(dict(
        sites=sites, price=price, gpr=gpr, egi=egi, opex=opex, noi=noi,
        oer=(opex / egi if egi else 0), caprate=(noi / price if price else 0),
        ads=ads, cfbt=noi - ads, coc=((noi - ads) / cash if cash else 0),
        dscr=(noi / ads if ads else float("inf")), debt_yield=(noi / loan if loan else 0),
        grm=(price / egi if egi else 0), spread=(noi / price if price else 0) - intr,
        one_pct=((egi / 12) / price if price else 0),
        quick_dscr=(noi / (price * intr) if price and intr else 0),
        loan=loan, down=down, closing=closing, cash=cash, intr=intr, ltv=ltv,
        exitcap=exitcap, hold=hold, price_per_site=(price / sites if sites else 0),
        noi_per_site=(noi / sites if sites else 0),
    ))

    # ---- multi-year IRR / equity multiple ----
    cfs = [-cash]
    for y in range(1, 11):
        g_gpr = gpr * (1 + rent_g) ** (y - 1)
        e_egi = g_gpr * (1 - vac) + other_a * (1 + rent_g) ** (y - 1)
        e_noi = e_egi - opex * (1 + exp_g) ** (y - 1)
        a = ads if y <= hold else 0
        ocf = (e_noi - a) if y <= hold else 0
        sale = 0
        if y == hold:
            k = min(y, amort) * 12
            i = intr / 12
            bal = loan * (1 + i) ** k - pmt_m * (((1 + i) ** k - 1) / i) if i else 0
            sale = (e_noi * (1 + rent_g)) / exitcap * (1 - sellcost) - bal if exitcap else 0
        cfs.append(ocf + sale)
    d["irr"] = irr(cfs)
    inflow = sum(c for c in cfs[1:])
    d["equity_multiple"] = (inflow / cash) if cash else 0

    # ---- normalization (seller vs normalized) ----
    if "Normalization" in wb.sheetnames:
        nm = wb["Normalization"]
        mgp = gv(nm, "B5", 0.10); rmp = gv(nm, "B6", 0.05)
        cxp = gv(nm, "B7", 0.03); txb = gv(nm, "B8", 0.20)
        n_egi = gv(nm, "B11")
        s_tax = gv(nm, "B15"); s_ins = gv(nm, "B16"); s_util = gv(nm, "B17")
        s_rm = gv(nm, "B18"); s_mgmt = gv(nm, "B19"); s_capex = gv(nm, "B20")
        s_pay = gv(nm, "B21"); s_admin = gv(nm, "B22")
        s_opex = s_tax + s_ins + s_util + s_rm + s_mgmt + s_capex + s_pay + s_admin
        n_opex = (s_tax * (1 + txb) + s_ins + s_util + max(s_rm, rmp * n_egi)
                  + max(s_mgmt, mgp * n_egi) + max(s_capex, cxp * n_egi) + s_pay + s_admin)
        valcap = gv(nm, "B30", 0.10)
        d["norm"] = dict(
            egi=n_egi, seller_opex=s_opex, norm_opex=n_opex,
            seller_oer=(s_opex / n_egi if n_egi else 0), norm_oer=(n_opex / n_egi if n_egi else 0),
            seller_noi=n_egi - s_opex, norm_noi=n_egi - n_opex,
            haircut=(n_egi - s_opex) - (n_egi - n_opex),
            overpay=(((n_egi - s_opex) - (n_egi - n_opex)) / valcap if valcap else 0),
        )

    # ---- offer structures ----
    if "Offer Structures" in wb.sheetnames:
        of = wb["Offer Structures"]
        tdscr = gv(of, "B5", 1.35); tcoc = gv(of, "B6", 0.10)
        # MAO uses the Offer tab's own levers (falling back to Underwriting),
        # so the report matches what the workbook's MAO block shows
        m_ltv = gv(of, "B7", 0) or ltv
        m_rate = gv(of, "B8", 0) or intr
        m_amort = gv(of, "B9", 0) or amort
        k = pmt(m_rate / 12, m_amort * 12, -1) * 12
        mao_dscr = noi / (tdscr * m_ltv * k) if (tdscr and m_ltv and k) else 0
        mao_coc = ((noi - tcoc * capex) / (m_ltv * k + tcoc * (1 - m_ltv + close_pct))) if (m_ltv * k + tcoc * (1 - m_ltv + close_pct)) else 0
        d["mao"] = max(0, min(mao_dscr, mao_coc))
        d["tdscr"], d["tcoc"] = tdscr, tcoc
        offers = []
        names = {"B": "Seller Asking", "C": "Conventional", "D": "Partial Carry", "E": "Full Carry"}
        for col in "BCDE":
            p = gv(of, f"{col}18")
            bp = gv(of, f"{col}19"); br = gv(of, f"{col}20"); ba = gv(of, f"{col}21")
            bio = str(of[f"{col}22"].value or "N").strip().upper()
            cp = gv(of, f"{col}23"); cr = gv(of, f"{col}24"); ca = gv(of, f"{col}25")
            cio = str(of[f"{col}26"].value or "N").strip().upper()
            bankloan = p * bp
            bankads = bankloan * br if bio == "Y" else (pmt(br / 12, ba * 12, -bankloan) * 12 if bankloan else 0)
            carry = p * cp
            carryads = 0 if carry == 0 else (carry * cr if cio == "Y" else pmt(cr / 12, ca * 12, -carry) * 12)
            totads = bankads + carryads
            dn = p * (1 - bp - cp); clcap = p * close_pct + capex; csh = dn + clcap
            dscr = (noi / totads) if totads else float("inf")
            cfbt = noi - totads; coc = (cfbt / csh) if csh else 0
            dpass = (totads == 0) or (dscr >= tdscr); cpass = coc >= tcoc
            verdict = "GO" if (dpass and cpass) else ("CONDITIONAL" if (dpass or cpass) else "NO-GO")
            offers.append(dict(name=names[col], price=p, dscr=dscr, coc=coc,
                               cash=csh, cfbt=cfbt, verdict=verdict, blended=(
                                   (bankloan * br + carry * cr) / (bankloan + carry) if (bankloan + carry) else 0)))
        d["offers"] = offers

    # ---- scorecard verdict + red flags (uses the workbook's editable targets) ----
    v_dscr = d.get("tdscr", 1.35); v_coc = d.get("tcoc", 0.10)
    d["verdict"] = ("GO" if (d["dscr"] >= v_dscr and d["coc"] >= v_coc)
                    else ("CONDITIONAL" if (d["dscr"] >= v_dscr or d["coc"] >= v_coc) else "NO-GO"))
    flags = []
    if "norm" in d and d["norm"]["seller_oer"] and d["norm"]["seller_oer"] < 0.30:
        flags.append(f"Seller expense ratio {pct(d['norm']['seller_oer'])} is under 30% — NOI likely inflated (Linda: 'fiction').")
    if d["dscr"] < v_dscr:
        flags.append(f"DSCR {mult(d['dscr'])} is below the {mult(v_dscr)} floor — financing risk.")
    if d["coc"] < v_coc:
        flags.append(f"Cash-on-cash {pct(d['coc'])} is below the {pct(v_coc,0)} target.")
    if d["caprate"] < d["intr"]:
        flags.append(f"Cap rate {pct(d['caprate'])} is below interest {pct(d['intr'])} — negative leverage (losing money day one).")
    if d["grm"] > 10:
        flags.append(f"GRM {mult(d['grm'])} is over 10 — priced high vs. income.")
    if "norm" in d and d["norm"]["haircut"] > 0:
        flags.append(f"Normalizing the NOI cuts {money(d['norm']['haircut'])} — ~{money(d['norm']['overpay'])} of overpayment risk.")
    d["flags"] = flags or ["No major red flags on the screening metrics — still verify the T-12 and condition."]
    return d


VALUE_ADD = [
    "Raise below-market lot rents to market (often the biggest lever).",
    "Improve occupancy — usually a management problem, not a real estate one.",
    "Convert owner-paid utilities to ratio bill-back (moves cost off the P&L).",
    "Add revenue: laundry, storage, covered parking, camp store.",
    "Don't pay the seller for value YOU will create — that upside is your return.",
]


# --------------------------------------------------------------------------- #
# PDF (one page)
# --------------------------------------------------------------------------- #
def draw_brandmark(c, x, ytop, h, BRAND, BLACK):
    """Vector Wise Certified emblem — mountain range, bold blue roofline, house+chimney,
    and a window — a faithful rendition of the logo. Returns width drawn."""
    w = h * 1.30
    yb = ytop - h
    # mountain range (black) — upper portion, two main peaks
    p = c.beginPath()
    p.moveTo(x, yb + 0.48 * h)
    p.lineTo(x + 0.18 * w, yb + 0.86 * h)
    p.lineTo(x + 0.30 * w, yb + 0.66 * h)
    p.lineTo(x + 0.50 * w, yb + 1.00 * h)
    p.lineTo(x + 0.66 * w, yb + 0.70 * h)
    p.lineTo(x + 0.82 * w, yb + 0.90 * h)
    p.lineTo(x + w, yb + 0.52 * h)
    p.lineTo(x + w, yb + 0.48 * h)
    p.close()
    c.setFillColor(BLACK); c.drawPath(p, fill=1, stroke=0)
    apex_x = x + 0.43 * w
    # small house + chimney sitting on the right slope of the roof
    hx = x + 0.66 * w; hw = 0.15 * w; hy = yb + 0.30 * h; hh = 0.17 * h
    c.setFillColor(BLACK)
    c.rect(hx, hy, hw, hh, fill=1, stroke=0)                       # body
    tp = c.beginPath(); tp.moveTo(hx - 0.025 * w, hy + hh)
    tp.lineTo(hx + hw / 2, hy + hh + 0.11 * h); tp.lineTo(hx + hw + 0.025 * w, hy + hh); tp.close()
    c.drawPath(tp, fill=1, stroke=0)                              # roof
    c.rect(hx + hw * 0.66, hy + hh + 0.02 * h, 0.03 * w, 0.11 * h, fill=1, stroke=0)  # chimney
    # bold royal-blue roofline chevron (front)
    c.setLineCap(1); c.setLineJoin(1)
    c.setStrokeColor(BRAND); c.setLineWidth(max(3.2, 0.15 * h))
    c.line(x + 0.04 * w, yb + 0.18 * h, apex_x, yb + 0.56 * h)
    c.line(apex_x, yb + 0.56 * h, x + 0.99 * w, yb + 0.12 * h)
    # thin dark second roofline (behind, for the layered look)
    c.setStrokeColor(BLACK); c.setLineWidth(max(1.1, 0.04 * h))
    c.line(x + 0.10 * w, yb + 0.12 * h, apex_x, yb + 0.46 * h)
    c.line(apex_x, yb + 0.46 * h, x + 0.92 * w, yb + 0.10 * h)
    # 2x2 window under the apex (black, on white)
    s = 0.052 * w; gap = 0.016 * w
    bx = apex_x - s - gap / 2; by = yb + 0.20 * h
    c.setFillColor(BLACK)
    for ix in (0, 1):
        for iy in (0, 1):
            c.rect(bx + ix * (s + gap), by + iy * (s + gap), s, s, fill=1, stroke=0)
    return w


def render_pdf(d, path, logo=None):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas

    BRAND = colors.HexColor("#1D3FAE")   # Wise Certified royal blue
    BLACK = colors.HexColor("#111111")
    BLUE = BRAND
    GREEN = colors.HexColor("#C6EFCE"); YELL = colors.HexColor("#FFEB9C"); RED = colors.HexColor("#FFC7CE")
    vfill = {"GO": GREEN, "CONDITIONAL": YELL, "NO-GO": RED}

    c = canvas.Canvas(path, pagesize=letter)
    W, H = letter
    m = 0.55 * inch
    top = H - m

    GREY = colors.HexColor("#555555")
    # ---- branded header: real logo if present, else vector Wise Certified mark ----
    logo_h = 1.00 * inch
    drew = False
    if logo and os.path.exists(logo):
        try:
            from reportlab.lib.utils import ImageReader
            img = ImageReader(logo); iw, ih = img.getSize(); ar = iw / ih if ih else 1
            c.drawImage(img, m, top - logo_h, width=logo_h * ar, height=logo_h,
                        mask="auto", preserveAspectRatio=True)
            drew = True
        except Exception:
            drew = False
    if not drew:
        iconw = draw_brandmark(c, m, top, logo_h, BRAND, BLACK)
        wx = m + iconw + 0.16 * inch
        c.setFillColor(BLACK); c.setFont("Helvetica-Bold", 20)
        c.drawString(wx, top - 0.36 * inch, "WISE CERTIFIED")
        c.setFillColor(BRAND); c.setFont("Helvetica-Bold", 9)
        c.drawString(wx, top - 0.56 * inch, "H O M E   B U Y E R S")

    # verdict badge (top right)
    c.setFillColor(vfill.get(d["verdict"], YELL))
    c.roundRect(W - m - 1.9 * inch, top - 0.58 * inch, 1.9 * inch, 0.55 * inch, 6, fill=1, stroke=0)
    c.setFillColor(colors.black); c.setFont("Helvetica-Bold", 15)
    c.drawCentredString(W - m - 0.95 * inch, top - 0.36 * inch, d["verdict"])

    # brand rule
    rule_y = top - logo_h - 0.06 * inch
    c.setStrokeColor(BRAND); c.setLineWidth(2); c.line(m, rule_y, W - m, rule_y)

    # deal title line under the rule
    c.setFillColor(BRAND); c.setFont("Helvetica-Bold", 15)
    c.drawString(m, rule_y - 0.26 * inch, str(d["name"])[:44])
    c.setFillColor(GREY); c.setFont("Helvetica", 10)
    c.drawString(m, rule_y - 0.44 * inch, f"{d['loc']}  ·  RV Park Underwriting Summary")
    y = rule_y - 0.72 * inch

    def header(label):
        nonlocal y
        c.setFillColor(BLUE); c.rect(m, y - 0.05 * inch, W - 2 * m, 0.24 * inch, fill=1, stroke=0)
        c.setFillColor(colors.white); c.setFont("Helvetica-Bold", 10)
        c.drawString(m + 4, y + 0.02 * inch, label.upper())
        y -= 0.34 * inch

    def kvrow(pairs):
        nonlocal y
        colw = (W - 2 * m) / len(pairs)
        for i, (k, v) in enumerate(pairs):
            x = m + i * colw
            c.setFillColor(colors.HexColor("#666666")); c.setFont("Helvetica", 8)
            c.drawString(x, y + 0.14 * inch, k)
            c.setFillColor(colors.black); c.setFont("Helvetica-Bold", 13)
            c.drawString(x, y - 0.02 * inch, v)
        y -= 0.5 * inch

    header("5 Key Metrics")
    kvrow([("NOI", money(d["noi"])), ("Cap Rate", pct(d["caprate"])),
           ("DSCR", mult(d["dscr"])), ("Cash Flow", money(d["cfbt"])), ("Cash-on-Cash", pct(d["coc"]))])

    header("Property Snapshot")
    kvrow([("Price", money(d["price"])), ("Sites", f"{int(d['sites'])}"),
           ("Price/Site", money(d["price_per_site"])), ("Debt Yield", pct(d["debt_yield"])),
           ("5-yr IRR", pct(d["irr"]) if d.get("irr") is not None else "—")])

    if "norm" in d:
        header("Normalized vs Seller (build your own NOI)")
        n = d["norm"]
        kvrow([("Seller NOI", money(n["seller_noi"])), ("Normalized NOI", money(n["norm_noi"])),
               ("Seller Exp%", pct(n["seller_oer"])), ("NOI Haircut", money(n["haircut"])),
               ("Overpay Risk", money(n["overpay"]))])

    if "offers" in d:
        header(f"Offer Structures   (MAO {money(d['mao'])})")
        c.setFont("Helvetica-Bold", 8.5); c.setFillColor(colors.HexColor("#666666"))
        cols = [m, m + 1.9 * inch, m + 3.0 * inch, m + 4.0 * inch, m + 5.2 * inch]
        for x, t in zip(cols, ["Structure", "Price", "DSCR", "Cash-on-Cash", "Verdict"]):
            c.drawString(x, y + 0.12 * inch, t)
        y -= 0.06 * inch
        for o in d["offers"]:
            y -= 0.235 * inch
            c.setFillColor(vfill.get(o["verdict"], YELL))
            c.roundRect(cols[4], y - 0.02 * inch, 1.2 * inch, 0.2 * inch, 3, fill=1, stroke=0)
            c.setFillColor(colors.black); c.setFont("Helvetica", 9.5)
            c.drawString(cols[0], y, o["name"])
            c.drawString(cols[1], y, money(o["price"]))
            c.drawString(cols[2], y, mult(o["dscr"]) if o["dscr"] != float("inf") else "n/a")
            c.drawString(cols[3], y, pct(o["coc"]))
            c.setFont("Helvetica-Bold", 9.5)
            c.drawCentredString(cols[4] + 0.6 * inch, y, o["verdict"])
        y -= 0.34 * inch

    header("Red Flags")
    c.setFillColor(colors.black); c.setFont("Helvetica", 9)
    for f in d["flags"][:5]:
        c.drawString(m + 4, y, u"• " + f); y -= 0.2 * inch
    y -= 0.1 * inch

    header("Value-Add Upside")
    c.setFillColor(colors.black); c.setFont("Helvetica", 9)
    for v in VALUE_ADD:
        c.drawString(m + 4, y, u"• " + v); y -= 0.2 * inch

    c.setFillColor(colors.HexColor("#999999")); c.setFont("Helvetica-Oblique", 7.5)
    c.drawString(m, 0.4 * inch, "© 2026 Wise Certified Home Buyers  ·  Screening tool, not investment advice  "
                               "·  Built with LindaAI.")
    c.showPage(); c.save()


# --------------------------------------------------------------------------- #
# PPTX (~6 slides)
# --------------------------------------------------------------------------- #
def render_pptx(d, path, logo=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY = RGBColor(0x1D, 0x3F, 0xAE)   # Wise Certified royal blue
    GREEN = RGBColor(0xC6, 0xEF, 0xCE); YELL = RGBColor(0xFF, 0xEB, 0x9C); RED = RGBColor(0xFF, 0xC7, 0xCE)
    vfill = {"GO": GREEN, "CONDITIONAL": YELL, "NO-GO": RED}
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide(title):
        s = prs.slides.add_slide(blank)
        bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        tf = bar.text_frame; tf.word_wrap = True
        tf.text = title; tf.paragraphs[0].font.size = Pt(30)
        tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        bar.text_frame.margin_left = Inches(0.4)
        return s

    def bullets(s, items, top=1.4, size=18):
        tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(5.6))
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = u"• " + it; p.font.size = Pt(size)

    # 1 — title / verdict
    s = slide(d["name"])
    if logo and os.path.exists(logo):
        try:
            s.shapes.add_picture(logo, Inches(0.6), Inches(1.4), height=Inches(1.6))
        except Exception:
            pass
    else:
        wm = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(6), Inches(1.0))
        wm.text_frame.text = "WISE CERTIFIED  ·  HOME BUYERS"
        wm.text_frame.paragraphs[0].font.size = Pt(20); wm.text_frame.paragraphs[0].font.bold = True
        wm.text_frame.paragraphs[0].font.color.rgb = NAVY
    bullets(s, [d["loc"], "", f"VERDICT: {d['verdict']}",
                f"Price {money(d['price'])}  ·  {int(d['sites'])} sites  ·  {money(d['price_per_site'])}/site"], top=3.2, size=22)
    badge = s.shapes.add_shape(1, Inches(9.6), Inches(1.4), Inches(3.0), Inches(1.1))
    badge.fill.solid(); badge.fill.fore_color.rgb = vfill.get(d["verdict"], YELL); badge.line.fill.background()
    badge.text_frame.text = d["verdict"]
    badge.text_frame.paragraphs[0].font.size = Pt(28); badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 2 — 5 key metrics
    s = slide("The 5 Key Metrics")
    bullets(s, [f"NOI:  {money(d['noi'])}   ({money(d['noi_per_site'])}/site)",
                f"Cap Rate:  {pct(d['caprate'])}",
                f"DSCR:  {mult(d['dscr'])}   (Linda floor 1.35x)",
                f"Net Profit (cash flow):  {money(d['cfbt'])}",
                f"Cash-on-Cash:  {pct(d['coc'])}   (Linda target 10%)",
                f"5-yr IRR:  {pct(d['irr']) if d.get('irr') is not None else '—'}   ·   Equity Multiple:  {mult(d['equity_multiple'])}"], size=20)

    # 3 — normalized vs seller
    if "norm" in d:
        n = d["norm"]; s = slide("Normalized NOI vs Seller's Numbers")
        bullets(s, [f"Seller NOI:  {money(n['seller_noi'])}   (expense ratio {pct(n['seller_oer'])})",
                    f"Normalized NOI:  {money(n['norm_noi'])}   (expense ratio {pct(n['norm_oer'])})",
                    f"NOI Haircut:  {money(n['haircut'])}",
                    f"Overpayment Risk:  {money(n['overpay'])}",
                    "",
                    "Linda's adjustments: management 10% · R&M 5% · capex 3% · taxes +20%."], size=20)

    # 4 — offer comparison
    if "offers" in d:
        s = slide(f"Offer Structures   (MAO {money(d['mao'])})")
        rows = [f"{o['name']}:  {money(o['price'])}  ·  DSCR {mult(o['dscr']) if o['dscr']!=float('inf') else 'n/a'}  ·  CoC {pct(o['coc'])}  ->  {o['verdict']}" for o in d["offers"]]
        bullets(s, rows, size=18)

    # 5 — risks / red flags
    s = slide("Risks & Red Flags")
    bullets(s, d["flags"], size=18)

    # 6 — value add
    s = slide("Value-Add Upside")
    bullets(s, VALUE_ADD, size=18)

    prs.save(path)


def main():
    ap = argparse.ArgumentParser(description="Generate PDF/PPTX reports from a filled underwriting workbook.")
    ap.add_argument("workbook")
    ap.add_argument("--pdf", action="store_true")
    ap.add_argument("--pptx", action="store_true")
    ap.add_argument("--out", default=None, help="output basename (default: workbook name)")
    default_logo = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "wise-certified-logo.png")
    ap.add_argument("--logo", default=default_logo, help="brand logo PNG (auto-embedded if present)")
    args = ap.parse_args()
    if not (args.pdf or args.pptx):
        args.pdf = args.pptx = True

    wb = openpyxl.load_workbook(args.workbook, data_only=False)
    if "Underwriting" not in wb.sheetnames:
        sys.exit("That workbook has no 'Underwriting' tab — is it the right file?")
    d = compute(wb)
    base = args.out or os.path.splitext(args.workbook)[0]
    made = []
    if args.pdf:
        p = base + ".pdf"; render_pdf(d, p, logo=args.logo); made.append(p)
    if args.pptx:
        p = base + ".pptx"; render_pptx(d, p, logo=args.logo); made.append(p)
    print("Generated: " + " , ".join(made))
    print(f"Verdict: {d['verdict']}  |  NOI {money(d['noi'])}  Cap {pct(d['caprate'])}  "
          f"DSCR {mult(d['dscr'])}  CoC {pct(d['coc'])}")


if __name__ == "__main__":
    main()
